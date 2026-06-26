import logging
from pathlib import Path
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)


def _save_slide_images(slide, slide_num: int, images_dir: Path) -> List[Dict[str, str]]:
    """
    Extract and save images from all picture shapes on a slide.
    Returns list of dicts with 'path' and 'filename' for each saved image.
    """
    saved_images = []
    images_dir.mkdir(parents=True, exist_ok=True)
    img_idx = 0

    for shape in slide.shapes:
        try:
            # Check if shape has an image (Picture shape type = 13)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext.lower() if image.ext else "png"

                if not image_bytes or len(image_bytes) < 500:
                    continue  # Skip tiny/corrupt

                # Normalise to PNG
                try:
                    from PIL import Image
                    import io
                    img = Image.open(io.BytesIO(image_bytes))
                    if img.mode not in ("RGB", "RGBA", "L"):
                        img = img.convert("RGB")
                    out_buf = io.BytesIO()
                    img.save(out_buf, format="PNG")
                    image_bytes = out_buf.getvalue()
                    image_ext = "png"
                except Exception:
                    pass  # Keep original bytes

                img_filename = f"slide_{slide_num}_img_{img_idx + 1}.{image_ext}"
                img_path = images_dir / img_filename

                with open(img_path, "wb") as f:
                    f.write(image_bytes)

                saved_images.append({
                    "path": str(img_path),
                    "filename": img_filename,
                    "page": slide_num
                })
                img_idx += 1
                logger.debug(f"Saved PPTX image: {img_filename}")

        except Exception as e:
            logger.warning(f"Failed to extract image from slide {slide_num}, shape '{getattr(shape, 'name', '?')}': {e}")

    return saved_images


def extract_pptx_data(file_path: Path, images_dir: Optional[Path] = None) -> dict:
    """
    Extracts text and images slide-by-slide from a PPTX file.

    Args:
        file_path: Path to the PPTX file.
        images_dir: Optional directory to save extracted slide images.

    Returns:
        dict: {
            "filename": str,
            "total_pages": int,
            "pages": [
                {
                    "page_number": int,
                    "text": str,
                    "images": [{"path": str, "filename": str, "page": int}]
                }
            ]
        }
    """
    from pptx import Presentation

    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    logger.info(f"Starting slide text extraction for: {file_path.name}")

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        pages_data = []

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            slide_text_runs = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_runs.append(shape.text.strip())

            # Extract text from notes if any
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text_runs.append(f"[Presenter Notes]\n{notes_text}")

            # Clean up whitespace and empty lines
            lines = [line.strip() for line in "\n".join(slide_text_runs).split("\n")]
            cleaned_text = "\n".join([line for line in lines if line])

            page_data: Dict[str, Any] = {
                "page_number": slide_num,
                "text": cleaned_text,
                "images": []
            }

            # Extract images from slide if images_dir is provided
            if images_dir:
                saved = _save_slide_images(slide, slide_num, images_dir)
                page_data["images"] = saved

            pages_data.append(page_data)

        logger.info(f"Successfully extracted {total_slides} slides from {file_path.name}")
        return {
            "filename": file_path.name,
            "total_pages": total_slides,
            "pages": pages_data
        }
    except Exception as e:
        logger.error(f"Error reading PPTX {file_path.name}: {e}")
        raise e
