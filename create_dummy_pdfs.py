import os
import sys
from pathlib import Path

# Try importing pptx and docx first
try:
    from docx import Document
    from pptx import Presentation
except ImportError:
    print("\n[!] Error: python-docx or python-pptx is not installed.")
    print("    Please run: pip install -r backend/requirements.txt")
    print("    Then try running this script again.\n")
    sys.exit(1)

def create_pdf(filepath, pages_text):
    """
    Generates a valid, multi-page PDF in pure Python with page-by-page text.
    """
    num_pages = len(pages_text)
    kids_refs = [f"{4 + 2*i} 0 R" for i in range(num_pages)]
    kids_str = " ".join(kids_refs)
    
    catalog_str = "1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
    pages_container_str = f"2 0 obj\n<< /Type /Pages /Kids [{kids_str}] /Count {num_pages} >>\nendobj\n"
    font_str = "3 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
    
    body = [catalog_str, pages_container_str, font_str]
    
    for i, page_text in enumerate(pages_text):
        escaped_text = page_text.replace("(", "\\(").replace(")", "\\)")
        lines = [line.strip() for line in escaped_text.split("\n") if line.strip()]
        
        # Build text stream commands
        stream_cmds = ["BT", "/F1 12 Tf", "50 780 Td", "16 TL"]
        for line in lines:
            # Word wrapping for PDF stream
            words = line.split(" ")
            line_chunk = []
            for word in words:
                line_chunk.append(word)
                if len(" ".join(line_chunk)) > 70:
                    stream_cmds.append(f"({' '.join(line_chunk[:-1])}) Tj T*")
                    line_chunk = [word]
            if line_chunk:
                stream_cmds.append(f"({' '.join(line_chunk)}) Tj T*")
        stream_cmds.append("ET")
        stream_content = "\n".join(stream_cmds)
        
        page_idx = 4 + 2*i
        content_idx = 4 + 2*i + 1
        
        page_obj = f"{page_idx} 0 obj\n<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 3 0 R >> >> /MediaBox [0 0 595 842] /Contents {content_idx} 0 R >>\nendobj\n"
        content_obj = f"{content_idx} 0 obj\n<< /Length {len(stream_content)} >>\nstream\n{stream_content}\nendstream\nendobj\n"
        
        body.append(page_obj)
        body.append(content_obj)
        
    offsets = {}
    pdf_bytes = bytearray("%PDF-1.4\n".encode("latin-1"))
    
    for idx, obj_str in enumerate(body[:3]):
        obj_id = idx + 1
        offsets[obj_id] = len(pdf_bytes)
        pdf_bytes.extend(obj_str.encode("latin-1"))
        
    for i in range(num_pages):
        page_id = 4 + 2*i
        content_id = 4 + 2*i + 1
        
        offsets[page_id] = len(pdf_bytes)
        pdf_bytes.extend(body[3 + 2*i].encode("latin-1"))
        
        offsets[content_id] = len(pdf_bytes)
        pdf_bytes.extend(body[3 + 2*i + 1].encode("latin-1"))
        
    startxref = len(pdf_bytes)
    total_objects = 3 + 2*num_pages
    
    xref_lines = [f"xref\n0 {total_objects + 1}\n0000000000 65535 f \n"]
    for obj_id in range(1, total_objects + 1):
        offset = offsets[obj_id]
        xref_lines.append(f"{offset:010d} 00000 n \n")
        
    xref_str = "".join(xref_lines)
    pdf_bytes.extend(xref_str.encode("latin-1"))
    
    trailer_str = f"trailer\n<< /Size {total_objects + 1} /Root 1 0 R >>\nstartxref\n{startxref}\n%%EOF\n"
    pdf_bytes.extend(trailer_str.encode("latin-1"))
    
    filepath.parent.mkdir(parents=True, exist_ok=True)
    with open(filepath, "wb") as f:
        f.write(pdf_bytes)
    print(f"Created PDF: {filepath} ({len(pdf_bytes)} bytes) with {num_pages} pages")

def create_pptx(filepath, slides_text):
    """
    Generates a valid PPTX slide deck with slide-by-slide text.
    """
    prs = Presentation()
    
    # slide_layouts[5] is a blank layout with just a slide title or completely blank
    # Let's use layout 1 (Title + Content) for slides
    slide_layout = prs.slide_layouts[1]
    
    for idx, slide_content in enumerate(slides_text):
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        lines = slide_content.strip().split("\n")
        title_shape.text = lines[0] if lines else f"Slide {idx + 1}"
        
        if len(lines) > 1:
            tf = body_shape.text_frame
            tf.text = lines[1]
            for extra_line in lines[2:]:
                p = tf.add_paragraph()
                p.text = extra_line
                
    filepath.parent.mkdir(parents=True, exist_ok=True)
    prs.save(filepath)
    print(f"Created PPTX: {filepath} with {len(slides_text)} slides")

def create_docx(filepath, pages_text):
    """
    Generates a DOCX document with text and explicit page breaks.
    """
    doc = Document()
    
    for idx, page_content in enumerate(pages_text):
        if idx > 0:
            doc.add_page_break()
            
        lines = page_content.strip().split("\n")
        doc.add_heading(lines[0] if lines else f"Section {idx + 1}", level=1)
        
        for extra_line in lines[1:]:
            doc.add_paragraph(extra_line)
            
    filepath.parent.mkdir(parents=True, exist_ok=True)
    doc.save(filepath)
    print(f"Created DOCX: {filepath} with {len(pages_text)} virtual pages")

# --- Sample Datasets ---
dbms_data = [
    "DBMS Notes - Introduction to Databases\nA Database Management System (DBMS) is software designed to store, retrieve, and manage data. It replaces flat file systems by centralizing storage. Main benefits: consistency, integrity, security.",
    "DBMS Notes - Relational Model & Keys\nIn the Relational Model, tables represent relations. Rows represent tuples, columns represent attributes. A primary key uniquely identifies rows. Candidate keys are minimal unique keys. Foreign keys reference primary keys.",
    "DBMS Notes - Normalization Theory\nNormalization reduces data redundancy and prevents anomalies. 1NF: Atomic values. 2NF: In 1NF and no partial key dependencies. 3NF: In 2NF and no transitive dependencies. BCNF: If A -> B, A must be a super key."
]

os_slides = [
    "OS Lecture - Introduction to Operating Systems\nWhat is an OS?\n- Intermediate program between hardware and user.\n- Goals: resource efficiency, ease of execution.\n- Kernel: Heart of OS, always in memory.",
    "OS Lecture - Process Management\nProcess States:\n- New: Process is being created.\n- Ready: Waiting to be assigned to processor.\n- Running: Instructions are executing.\n- Waiting: Blocked on event or I/O.\n- Terminated: Execution finished.",
    "OS Lecture - Deadlocks & Critical Section\nDeadlock definition:\n- A set of processes blocked waiting for resources.\n- 4 Conditions: Mutual exclusion, hold & wait, no preemption, circular wait.\n- Solution: Banker's algorithm safe state check."
]

cn_doc = [
    "Computer Networks - Protocol Stack\nOSI 7-Layer Reference Model\n- Physical Layer: Transmission of raw bit streams.\n- Data Link Layer: Node-to-node framing and error detection.\n- Network Layer: Routing of packets (IP routing).\n- Transport Layer: End-to-end communication (TCP/UDP flow control).",
    "Computer Networks - Transport Layer Protocols\nTCP vs UDP comparison:\n- TCP is connection-oriented, reliable, and uses sliding window.\n- UDP is connectionless, fast, and does not perform retransmissions.\n- Three-Way Handshake establishes TCP: SYN, SYN-ACK, ACK."
]

oop_data = [
    "OOP Notes - Programming Fundamentals\nObject-Oriented Programming (OOP) uses classes and objects. A class is a template; an object is an instance. State, behavior, and identity are core characteristics.",
    "OOP Notes - Encapsulation & Polymorphism\nEncapsulation binds data and code together. Access modifiers: public, private, protected. Polymorphism allows method overloading (compile time) and method overriding (runtime)."
]

def main():
    notes_dir = Path("notes")
    
    # 1. Create a PDF under DBMS subfolder
    create_pdf(notes_dir / "DBMS" / "relational_theory.pdf", dbms_data)
    
    # 2. Create a PPTX under Operating Systems subfolder
    create_pptx(notes_dir / "Operating Systems" / "cpu_scheduling.pptx", os_slides)
    
    # 3. Create a DOCX under Computer Networks subfolder
    create_docx(notes_dir / "Computer Networks" / "osi_reference.docx", cn_doc)
    
    # 4. Create a PDF under OOP subfolder
    create_pdf(notes_dir / "OOP" / "oop_basics.pdf", oop_data)
    
    print("\nAll dummy notes folders and multi-format files created successfully!")

if __name__ == "__main__":
    main()
